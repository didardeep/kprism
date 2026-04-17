from pptx import Presentation
from pptx.util import Inches


def extract_text_from_ppt(file_path):
    """Extract text from all slides in a PPT file."""
    prs = Presentation(file_path)
    slides_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []

        for shape in slide.shapes:
            # Extract from text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            # Extract from tables
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        slide_texts.append(" | ".join(row_text))

        if slide_texts:
            slides_text.append({
                "slide_number": slide_num,
                "text": "\n".join(slide_texts)
            })

    return slides_text


def chunk_text(slides_text, chunk_size=500, chunk_overlap=50):
    """
    Break slide text into overlapping chunks.
    chunk_size and chunk_overlap are in words (approximate tokens).
    """
    chunks = []

    for slide in slides_text:
        slide_num = slide["slide_number"]
        text = slide["text"]
        words = text.split()

        if len(words) <= chunk_size:
            # Slide fits in one chunk
            chunks.append({
                "slide_number": slide_num,
                "chunk_text": text
            })
        else:
            # Split into overlapping chunks
            start = 0
            while start < len(words):
                end = start + chunk_size
                chunk_words = words[start:end]
                chunk_text = " ".join(chunk_words)
                chunks.append({
                    "slide_number": slide_num,
                    "chunk_text": chunk_text
                })
                start += chunk_size - chunk_overlap

    return chunks


def process_ppt(file_path):
    """Full pipeline: extract text from PPT and chunk it."""
    slides_text = extract_text_from_ppt(file_path)
    if not slides_text:
        return []
    chunks = chunk_text(slides_text)
    return chunks
