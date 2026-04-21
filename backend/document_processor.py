import os
from pptx import Presentation
from azure.ai.documentintelligence import DocumentIntelligenceClient
from azure.core.credentials import AzureKeyCredential
from config import AZURE_DI_ENDPOINT, AZURE_DI_KEY


def extract_text_from_ppt(file_path):
    """Extract text from all slides in a PPT/PPTX file."""
    prs = Presentation(file_path)
    pages_text = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        slide_texts.append(" | ".join(row_text))

        if slide_texts:
            pages_text.append({
                "page_number": slide_num,
                "text": "\n".join(slide_texts)
            })

    return pages_text


def extract_text_from_pdf(file_path):
    """Extract text from all pages in a PDF using Azure Document Intelligence."""
    client = DocumentIntelligenceClient(
        endpoint=AZURE_DI_ENDPOINT,
        credential=AzureKeyCredential(AZURE_DI_KEY),
    )

    with open(file_path, "rb") as f:
        poller = client.begin_analyze_document(
            "prebuilt-layout",
            analyze_request=f,
            content_type="application/octet-stream",
            output_content_format="markdown",
        )
    result = poller.result()

    # Group content by page
    pages_text = []
    if result.pages:
        for page in result.pages:
            page_num = page.page_number
            # Collect all lines on this page
            lines = []
            if page.lines:
                for line in page.lines:
                    text = line.content.strip()
                    if text:
                        lines.append(text)
            if lines:
                pages_text.append({
                    "page_number": page_num,
                    "text": "\n".join(lines),
                })

    return pages_text


def extract_text(file_path):
    """Extract text from a file based on its extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext in (".pptx", ".ppt"):
        return extract_text_from_ppt(file_path)
    elif ext == ".pdf":
        return extract_text_from_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def chunk_text(pages_text, chunk_size=500, chunk_overlap=50):
    """
    Break page text into overlapping chunks.
    chunk_size and chunk_overlap are in words (approximate tokens).
    """
    chunks = []

    for page in pages_text:
        page_num = page["page_number"]
        text = page["text"]
        words = text.split()

        if len(words) <= chunk_size:
            chunks.append({
                "page_number": page_num,
                "chunk_text": text
            })
        else:
            start = 0
            while start < len(words):
                end = start + chunk_size
                chunk_words = words[start:end]
                chunk_text = " ".join(chunk_words)
                chunks.append({
                    "page_number": page_num,
                    "chunk_text": chunk_text
                })
                start += chunk_size - chunk_overlap

    return chunks


def process_document(file_path):
    """Full pipeline: extract text from document and chunk it."""
    pages_text = extract_text(file_path)
    if not pages_text:
        return []
    chunks = chunk_text(pages_text)
    return chunks
